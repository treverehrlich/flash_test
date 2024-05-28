import os
from typing import Union
import pandas as pd
from docx import Document
from pptx import Presentation
from typing import List
from io import StringIO, BytesIO
from davinci.services.auth import get_s3_client, get_secret
from davinci.utils.logging import log, logger
from davinci.utils.df_engines import _get_engine, _get_read_func, _get_save_func, _VALID_ENGINE_TYPE
from davinci.utils.fileio import force_folder_to_path

@log()
def file_exists(path: str) -> bool:
    """
    Check if an S3 file exists.

    :param path: the path to the file on s3. Do not include the
        top level bucket name.
    :type s3_path: str
    :return: Boolean

    Example usage:

        .. code-block:: python

            if file_exists('prod/project/file.txt'):
                print('success')
    """
    s3 = get_s3_client()
    bucket = get_secret('AWS_BUCKET_NAME')
    try:
        s3.get_object(
            Bucket=bucket,
            Key=path
        )
        return True
    except s3.exceptions.NoSuchKey:
        return False


@log()
def get_file(s3_path: str, local_path: str, use_cache: bool=False) -> str:
    """
    Download a file from s3 to local machine.

    :param s3_path: the path to the file on s3. Do not include the
        top level bucket name.
    :type s3_path: str
    :param local_path: The local path to save the file to. Include the filename.
    :type local_path: str
    :param use_cache: Attempt to load the file from local machine before querying S3.
    :type use_cache: bool
    :return: the local path argument

    Example usage:

        .. code-block:: python

            s3_path('prod/project/file.txt', '/home/ubuntu/file.txt')
    """
    bucket_name = get_secret('AWS_BUCKET_NAME')
    force_folder_to_path(local_path)
    if use_cache and os.path.exists(local_path):
        return local_path
    s3 = get_s3_client()
    s3.download_file(bucket_name, s3_path, local_path)
    return local_path

@log()
def upload_file(local_path: str, s3_path: str, **kwargs) -> bool:
    """
    Upload a file to the Kenco SCS S3 bucket.

    :param local_path: local path to the file to upload
    :type local_path: str
    :param s3_path: Full S3 path (filename should be included)
    :type s3_path: str
    :return: None

    Example usage:

        .. code-block:: python

            upload_file('my/local/path/my_data.txt', 'my/s3/path/my_data.txt')
    """

    bucket_name = get_secret('AWS_BUCKET_NAME')
    s3 = get_s3_client()
    try:
        s3.upload_file(local_path, bucket_name, s3_path, **kwargs)
        return True
    except:
        logger.error(f'Error. Could not upload {local_path} to s3.')
        return False

@log()
def delete_file(s3_path: str) -> bool:
    """
    Upload a file to the Kenco SCS S3 bucket.

    :param local_path: local path to the file to upload
    :type local_path: str
    :param s3_path: Full S3 path (filename should be included)
    :type s3_path: str
    :return: None

    Example usage:

        .. code-block:: python

            upload_file('my/local/path/my_data.txt', 'my/s3/path/my_data.txt')
    """

    bucket_name = get_secret('AWS_BUCKET_NAME')
    s3 = get_s3_client()
    try:
        s3.delete_object(Bucket=bucket_name, Key=s3_path)
        return True
    except Exception as e:
        logger.error(f'Error. Could not delete {s3_path} from s3.')
        return False


@log()
def upload_df(df: pd.DataFrame, path: str, file_type: str ='csv') -> None:
    """
    Upload a DataFrame to the Kenco SCS S3 bucket.

    :param df: DataFrame to upload
    :type df: pd.DataFrame
    :param path: Full S3 path (filename should be included)
    :type path: str
    :param file_type: Can be 'csv', 'excel', or 'parquet'
    :type file_type: str
    :return: None

    Example usage:

        .. code-block:: python

            upload_df(df, 'my/s3/path/my_data.xlsx', file_type='excel')
    """

    local_file_path = os.path.basename(path)
    engine = _get_engine(file_type)
    save_func = _get_save_func(df, file_type)
    if file_type == 'csv':
        save_func(local_file_path, header=True, index=False)
    elif file_type == 'excel':
        save_func(local_file_path, header=True, index=False, engine=engine)
    elif file_type == 'parquet':
        save_func(local_file_path, index=False, engine=engine)
    else:
        raise ValueError('Bad file_type upload for upload_df()')
    
    upload_file(local_file_path, path)
    os.unlink(local_file_path)
    return True

@log()
def download_to_df(path: str, file_type: str='csv', save_path: Union[None, str]=None, **kwargs) -> pd.DataFrame:
    """
    Download a DataFrame from the Kenco SCS S3 bucket.
    If save_path is specified, a local cache of the
    data will be made.

    :param path: Full S3 path (filename should be included)
    :type path: str
    :param file_type: Can be 'csv', 'excel', or 'parquet'
    :type file_type: str
    :param save_path: Local path (filename should be included)
    :type save_path: str
    :return: DataFrame

    Example usage:

        .. code-block:: python

            df = download_to_df('my/s3/path/my_data.xlsx', file_type='excel')
    """
    if file_type not in _VALID_ENGINE_TYPE:
        raise TypeError('Invalid file_type parsed')
    bucket_name = 'aws-scs-prod-bucket'
    s3 = get_s3_client()
    obj = s3.get_object(Bucket=bucket_name, Key=path)
    data = obj['Body'].read()

    read_func = _get_read_func(file_type)
    engine = _get_engine(file_type)
    df = read_func(BytesIO(data), engine=engine, **kwargs)
    if save_path:
        force_folder_to_path(save_path)
        save_func = _get_save_func(df, file_type)
        if file_type == 'csv':
            save_func(save_path, header=True, index=False)
        elif file_type == 'excel':
            save_func(save_path, header=True, index=False, engine=engine)
        elif file_type == 'parquet':
            save_func(save_path, index=False, engine=engine)
        else:
            raise ValueError('Bad file_type download for download_to_df()')
    return df

@log()
def read_s3_files_in_folder(s3_folder_path: str) -> List[str]:
    """
    Read the content of all files in a folder from S3 based on their types.

    :param s3_folder_path: the path to the folder on S3.
    :type s3_folder_path: str
    :return: a list of file contents as strings.
    """
    bucket_name = get_secret('AWS_BUCKET_NAME')
    s3 = get_s3_client()
    objects = s3.list_objects_v2(Bucket=bucket_name, Prefix=s3_folder_path)

    file_contents = []
    
    for obj in objects.get('Contents', []):
        s3_path = obj['Key']
        obj = s3.get_object(Bucket=bucket_name, Key=s3_path)
        file_content = obj['Body'].read()

        if s3_path.endswith(".txt"):
            file_contents.append(file_content.decode('utf-8'))
        elif s3_path.endswith(".docx"):
            doc = Document(BytesIO(file_content))
            file_contents.append("\n".join([par.text for par in doc.paragraphs]))
        elif s3_path.endswith(".pptx"):
            prs = Presentation(BytesIO(file_content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            file_contents.append(text)

    return file_contents

@log()
def get_file_names_in_folder(s3_folder_path:str) -> List[str]:
    """
    Read the content of all files in a folder from S3 based on their types.

    :param s3_folder_path: the path to the folder on S3.
    :type s3_folder_path: str
    :return: a list of file names as strings.
    """
    bucket_name = get_secret('AWS_BUCKET_NAME')
    s3 = get_s3_client()
    objects = s3.list_objects_v2(Bucket=bucket_name, Prefix=s3_folder_path)

    file_contents = []
    for obj in objects.get('Contents', []):
        file_contents.append(os.path.basename(obj['Key']).split('/')[-1])
    
    return file_contents


 

